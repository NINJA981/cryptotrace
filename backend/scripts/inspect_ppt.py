import pptx
import json

prs = pptx.Presentation('SIH2026 (1).pptx')
slide_data = []

for idx, slide in enumerate(prs.slides):
    shapes_info = []
    for s_idx, shape in enumerate(slide.shapes):
        shape_entry = {
            "id": shape.shape_id,
            "name": shape.name,
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height,
            "has_text": shape.has_text_frame,
            "has_table": shape.has_table,
            "text": []
        }
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                txt = p.text.strip()
                if txt:
                    shape_entry["text"].append(txt)
        if shape.has_table:
            table_rows = []
            for row in shape.table.rows:
                table_rows.append([cell.text.strip() for cell in row.cells])
            shape_entry["table"] = table_rows
        shapes_info.append(shape_entry)
    slide_data.append({
        "slide_index": idx + 1,
        "shapes": shapes_info
    })

with open("data/ppt_inspection.json", "w", encoding="utf-8") as f:
    json.dump(slide_data, f, indent=2, ensure_ascii=False)

print(f"Inspection complete. {len(prs.slides)} slides saved to data/ppt_inspection.json.")
