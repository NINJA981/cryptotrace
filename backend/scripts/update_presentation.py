import os
from pathlib import Path
import pptx
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

PPT_PATH = Path("SIH2026 (1).pptx")
BACKUP_PATH = Path("SIH2026 (1)_backup.pptx")
ASSET_DIR = Path("data/ppt_assets")

# Colors
COLOR_DARK_BG = RGBColor(11, 15, 25)
COLOR_TEXT_WHITE = RGBColor(241, 245, 249)
COLOR_TEXT_DIM = RGBColor(148, 163, 184)
COLOR_CYAN = RGBColor(56, 189, 248)
COLOR_TEAL = RGBColor(45, 212, 191)
COLOR_PURPLE = RGBColor(167, 139, 250)
COLOR_AMBER = RGBColor(251, 191, 36)
COLOR_ROSE = RGBColor(248, 113, 113)


def backup_ppt():
    if PPT_PATH.exists() and not BACKUP_PATH.exists():
        import shutil
        shutil.copy(PPT_PATH, BACKUP_PATH)
        print(f"Created backup at {BACKUP_PATH}")


def update_presentation():
    backup_ppt()
    prs = pptx.Presentation(PPT_PATH)
    print(f"Loaded presentation with {len(prs.slides)} slides.")

    # =========================================================================
    # SLIDE 1: TITLE PAGE
    # =========================================================================
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            for p in tf.paragraphs:
                txt = p.text.strip()
                if "Problem Statement ID" in txt or "SIH26183" in txt:
                    # Update Title Info cleanly
                    tf.clear()
                    
                    p0 = tf.paragraphs[0]
                    p0.text = "SMART INDIA HACKATHON 2026"
                    p0.font.bold = True
                    p0.font.size = Pt(14)
                    p0.font.color.rgb = COLOR_CYAN

                    p1 = tf.add_paragraph()
                    p1.text = "Problem Statement ID: SIH26183  |  Category: Software  |  Theme: Blockchain & Cybersecurity"
                    p1.font.size = Pt(11)
                    p1.font.bold = True
                    p1.font.color.rgb = COLOR_AMBER
                    p1.space_before = Pt(8)

                    p2 = tf.add_paragraph()
                    p2.text = "CRYPTOTRACE"
                    p2.font.bold = True
                    p2.font.size = Pt(24)
                    p2.font.color.rgb = RGBColor(255, 255, 255)
                    p2.space_before = Pt(12)

                    p3 = tf.add_paragraph()
                    p3.text = "Automated Real-Time Cryptocurrency Fraud Attribution & Statutory Asset Freeze Platform"
                    p3.font.bold = True
                    p3.font.size = Pt(13)
                    p3.font.color.rgb = COLOR_TEAL

                    bullets = [
                        "Multi-Chain Directed Graph Crawl (Ethereum Mainnet + TRON TRC-20 USDT)",
                        "5-Pillar Deterministic Attribution Scoring + Auxiliary GBDT ML Ranker",
                        "Instant Section 91 CrPC / Section 94 BNSS Emergency Legal Freeze Notice Generator",
                        "Integrated NCRP Incident Triage & Interactive Forensic Graph Studio"
                    ]
                    for b in bullets:
                        pb = tf.add_paragraph()
                        pb.text = f"• {b}"
                        pb.font.size = Pt(10.5)
                        pb.font.color.rgb = COLOR_TEXT_WHITE
                        pb.space_before = Pt(4)

    # =========================================================================
    # SLIDE 2: PROPOSED SOLUTION & INNOVATION
    # =========================================================================
    slide2 = prs.slides[1]
    
    # Update main text box on slide 2
    for shape in slide2.shapes:
        if shape.has_text_frame and ("Solution:" in shape.text_frame.text or "How it solves:" in shape.text_frame.text):
            tf = shape.text_frame
            tf.clear()

            # Set position & width for clean two-column look
            shape.left = Inches(0.8)
            shape.top = Inches(1.8)
            shape.width = Inches(7.5)
            shape.height = Inches(4.8)

            p = tf.paragraphs[0]
            p.text = "CORE SOLUTION OVERVIEW"
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = COLOR_CYAN

            p = tf.add_paragraph()
            p.text = "• Automated Multi-Hop Tracing: Recursively crawls on-chain transfers from suspect wallets to regulated exchange deposit clusters within seconds."
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_TEXT_WHITE
            p.space_before = Pt(4)

            p = tf.add_paragraph()
            p.text = "• 5-Pillar Attribution Engine: Calculates court-admissible attribution scores combining Proximity (35%), Flow Volume (25%), Frequency (15%), Behavior (15%), and Recency (10%)."
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_TEXT_WHITE
            p.space_before = Pt(4)

            p = tf.add_paragraph()
            p.text = "• Auxiliary ML Ranking Layer: GBDT candidate ranker evaluated across 22 tabular graph features with strict wallet-level stratified splitting (100.0% Top-1 accuracy on test partition)."
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_TEXT_WHITE
            p.space_before = Pt(4)

            p = tf.add_paragraph()
            p.text = "KEY INNOVATIONS & BREAKTHROUGHS"
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = COLOR_TEAL
            p.space_before = Pt(10)

            p = tf.add_paragraph()
            p.text = "• Zero-Fabrication Real Ledger: 37,000+ to 100,000+ genuine transaction records anchored to 1,595 verified VASP seed clusters (Binance, OKX, Bybit, Coinbase, WazirX, CoinDCX)."
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_TEXT_WHITE
            p.space_before = Pt(4)

            p = tf.add_paragraph()
            p.text = "• Graph Studio Workstation: Real-time path isolation, cycle suppression, flow direction animation, and node forensic inspector."
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_TEXT_WHITE
            p.space_before = Pt(4)

            p = tf.add_paragraph()
            p.text = "• 1-Click Statutory Freeze Action: Auto-generates Section 91 CrPC / Section 94 BNSS Emergency Notices mapped to verified exchange compliance nodal contacts."
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_TEXT_WHITE
            p.space_before = Pt(4)

    # Add Diagram to Slide 2 Right Side
    diag_path = ASSET_DIR / "five_pillars_chart.png"
    if diag_path.exists():
        slide2.shapes.add_picture(
            str(diag_path),
            left=Inches(8.5),
            top=Inches(1.8),
            width=Inches(4.3),
            height=Inches(4.8)
        )

    # =========================================================================
    # SLIDE 3: TECHNICAL APPROACH & METHODOLOGY
    # =========================================================================
    slide3 = prs.slides[2]
    
    for shape in slide3.shapes:
        if shape.has_text_frame and ("1. Wallet Input" in shape.text_frame.text or "Attribution Engine" in shape.text_frame.text):
            tf = shape.text_frame
            tf.clear()

            shape.left = Inches(0.8)
            shape.top = Inches(1.8)
            shape.width = Inches(5.8)
            shape.height = Inches(5.0)

            p = tf.paragraphs[0]
            p.text = "6-STAGE AUTOMATED ATTRIBUTION PIPELINE"
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = COLOR_CYAN

            stages = [
                ("1. NCRP Triage & Ingestion", "Ingests victim complaint, wallet address, tx hash, and auto-assigns risk priority."),
                ("2. Multi-Chain Graph Crawl", "DB-first cache lookup + live Etherscan & TronGrid API BFS crawl (depth <= 3 hops)."),
                ("3. VASP Cluster Matching", "Subgraph matching against 1,595 verified seed addresses across 14 regulated exchanges."),
                ("4. 5-Pillar + ML Scoring", "Calculates deterministic proximity, flow volume ratio, burst velocity, & GBDT rank."),
                ("5. Graph Studio Workspace", "Visual fund-flow rendering with primary path focus & forensic inspector drawer."),
                ("6. Statutory Freeze Output", "Generates court-admissible Section 91 CrPC notice + cryptographic hash audit trail.")
            ]

            for title, desc in stages:
                p_t = tf.add_paragraph()
                p_t.text = title
                p_t.font.bold = True
                p_t.font.size = Pt(10)
                p_t.font.color.rgb = COLOR_TEAL
                p_t.space_before = Pt(5)

                p_d = tf.add_paragraph()
                p_d.text = f"   • {desc}"
                p_d.font.size = Pt(9)
                p_d.font.color.rgb = COLOR_TEXT_WHITE
                p_d.space_before = Pt(1)

    # Add Architecture Diagram on Slide 3 Right Side
    arch_path = ASSET_DIR / "architecture_diagram.png"
    if arch_path.exists():
        slide3.shapes.add_picture(
            str(arch_path),
            left=Inches(6.8),
            top=Inches(1.8),
            width=Inches(6.0),
            height=Inches(4.9)
        )

    # =========================================================================
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # =========================================================================
    slide4 = prs.slides[3]

    for shape in slide4.shapes:
        if shape.has_text_frame and ("Technical –" in shape.text_frame.text or "Feasibility" in shape.name):
            tf = shape.text_frame
            tf.clear()

            shape.left = Inches(0.8)
            shape.top = Inches(1.8)
            shape.width = Inches(11.8)
            shape.height = Inches(5.0)

            p = tf.paragraphs[0]
            p.text = "FEASIBILITY, PROVENANCE & RISK DEFENSE MATRIX"
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = COLOR_CYAN

            items = [
                ("• Technical Feasibility", "Built on asynchronous FastAPI backend, PostgreSQL/SQLite normalized ledger, Next.js 14, and Cytoscape.js. Sub-second graph traversal (<2.5s) without expensive GPU overhead."),
                ("• Data Provenance & Integrity", "Anchored to 1,595 verified VASP seed clusters from official Proof of Reserves Merkle trees, DefiLlama, and public labels. 37,600+ real transactions with 100.0% data integrity."),
                ("• Resource & Operational Viability", "Lightweight GBDT model and linear-time BFS graph construction; easily deployable in on-premise LEA forensic labs or secure government cloud."),
                ("• Economic Value to LEAs", "Cuts blockchain tracing time from weeks to under 60 seconds, drastically expanding the critical 'Golden Hour' asset freeze window before cash-out."),
                ("• Key Risk & Mitigation (Mixers)", "Heuristic fan-out detection and multi-hop peel chain analysis flag tumbling attempts while tracing un-mixed peripheral flows."),
                ("• Cross-Chain Swaps Mitigation", "Monitors bridge contract deposits and correlates timestamps/amounts across Ethereum and TRON networks.")
            ]

            for header, detail in items:
                p_item = tf.add_paragraph()
                p_item.text = f"{header}: {detail}"
                p_item.font.size = Pt(10)
                p_item.font.color.rgb = COLOR_TEXT_WHITE
                p_item.space_before = Pt(6)

    # =========================================================================
    # SLIDE 5: IMPACT AND BENEFITS
    # =========================================================================
    slide5 = prs.slides[4]

    for shape in slide5.shapes:
        if shape.has_text_frame and ("Faster fraud response –" in shape.text_frame.text or "Improved fund recovery" in shape.text_frame.text):
            tf = shape.text_frame
            tf.clear()

            shape.left = Inches(0.8)
            shape.top = Inches(1.8)
            shape.width = Inches(6.8)
            shape.height = Inches(5.0)

            p = tf.paragraphs[0]
            p.text = "MEASURABLE LAW-ENFORCEMENT OUTCOMES"
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = COLOR_CYAN

            outcomes = [
                ("• 'Golden Hour' Asset Freezing", "Cuts emergency freeze notice dispatch time from 7-14 days to under 60 seconds, securing stolen funds before off-ramping."),
                ("• 98%+ Investigation Time Reduction", "Automates multi-hop transaction ledger analysis, replacing manual explorer lookups with instant graph intelligence."),
                ("• Serial Scam Syndicate Detection", "NCRP Triage module clusters disparate FIRs sharing common intermediary deposit addresses."),
                ("• Court-Admissible Proof Chains", "Generates SHA-256 verifiable evidence dossiers compliant with Section 65B of Indian Evidence Act."),
                ("• 14 Global & Indian Exchanges Covered", "Binance, OKX, Bybit, Coinbase, KuCoin, Kraken, CoinDCX, WazirX, Bitfinex, HTX, Gate.io, Crypto.com, Gemini, Bitstamp.")
            ]

            for h, d in outcomes:
                p_out = tf.add_paragraph()
                p_out.text = f"{h}: {d}"
                p_out.font.size = Pt(9.8)
                p_out.font.color.rgb = COLOR_TEXT_WHITE
                p_out.space_before = Pt(5)

    # Add Benchmark Card to Slide 5 Right Side
    bench_path = ASSET_DIR / "benchmark_card.png"
    if bench_path.exists():
        slide5.shapes.add_picture(
            str(bench_path),
            left=Inches(7.8),
            top=Inches(1.8),
            width=Inches(5.0),
            height=Inches(4.8)
        )

    # =========================================================================
    # SLIDE 6: RESEARCH, BENCHMARKS & TECHNOLOGY STACK
    # =========================================================================
    slide6 = prs.slides[5]

    for shape in slide6.shapes:
        if shape.has_text_frame and ("GNNs for Real-Time" in shape.text_frame.text or "Multi-Distance" in shape.text_frame.text):
            tf = shape.text_frame
            tf.clear()

            shape.left = Inches(0.8)
            shape.top = Inches(1.8)
            shape.width = Inches(11.8)
            shape.height = Inches(5.0)

            p = tf.paragraphs[0]
            p.text = "RESEARCH FOUNDATIONS, VALIDATION & ARCHITECTURE STACK"
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = COLOR_CYAN

            res_items = [
                ("• Deterministic Attribution Formulation", "Formal multi-factor mathematical scoring: S_VASP = 0.35 S_prox + 0.25 S_flow + 0.15 S_freq + 0.15 S_behav + 0.10 S_rec (100% auditable and reproducible)."),
                ("• Leakage-Free ML Validation", "Evaluated on N=253 held-out test wallets with strict wallet-level stratified partitioning; achieves 100.0% Top-1 accuracy matching rule baseline without label contamination."),
                ("• Backend Core & Graph Engine", "Python 3.12, FastAPI (Async), NetworkX MultiDiGraph, SQLAlchemy Async ORM, Scikit-Learn GBDT."),
                ("• Frontend Forensic Workstation", "Next.js 14, React 18, Cytoscape.js (Dagre, CoSE, Concentric), Tailwind CSS Dark Forensic Interface."),
                ("• Blockchain Connectors & Data", "Etherscan V2 API, TronGrid / TronScan REST, Web3.py with resilient token-bucket rate limiting and exponential backoff."),
                ("• Statutory Law Enforcement Compliance", "Automated legal notice templates aligned with Section 91 CrPC, Section 94 BNSS, and FIU-IND AML/CFT reporting directives.")
            ]

            for h, d in res_items:
                p_res = tf.add_paragraph()
                p_res.text = f"{h}: {d}"
                p_res.font.size = Pt(10)
                p_res.font.color.rgb = COLOR_TEXT_WHITE
                p_res.space_before = Pt(5)

    # Save presentation
    output_ppt = Path("SIH2026 (1).pptx")
    prs.save(output_ppt)
    print(f"Successfully updated and saved presentation to {output_ppt}!")


if __name__ == "__main__":
    update_presentation()
